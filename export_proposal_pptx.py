"""Generate ControlPlane.ai business proposal PowerPoint."""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x1E, 0x2F, 0x6E)
BLUE = RGBColor(0x3B, 0x5B, 0xDB)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF4, 0xFF)
PALE = RGBColor(0xF8, 0xFA, 0xFC)
LINE = RGBColor(0xCB, 0xD5, 0xE1)
RED = RGBColor(0xB9, 0x1C, 0x1C)
GREEN = RGBColor(0x04, 0x78, 0x57)
ORANGE = RGBColor(0xC2, 0x41, 0x0C)
VIOLET = RGBColor(0x6D, 0x28, 0xD9)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def rgb_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def box(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    rgb_fill(s, color)
    return s


def round_box(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    rgb_fill(s, color)
    s.adjustments[0] = 0.08
    return s


def set_run(run, text, size=14, bold=False, color=DARK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_text(slide, l, t, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size, bold, color, font)
    return tb


def add_para(tf, text, size=13, bold=False, color=DARK, space_before=4, space_after=2):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    set_run(p.add_run(), text, size, bold, color)
    return p


def footer(slide, page, total=20):
    box(slide, 0, Inches(7.28), W, Inches(0.22), NAVY)
    add_text(slide, Inches(0.4), Inches(7.28), Inches(8), Inches(0.22),
             "ControlPlane.ai  |  Accenture Innovation Challenge 2026  |  Round 2  |  Track 1",
             9, False, WHITE)
    add_text(slide, Inches(11.6), Inches(7.28), Inches(1.4), Inches(0.22),
             f"{page}  /  {total}", 9, False, WHITE, PP_ALIGN.RIGHT)


def header_bar(slide, title, subtitle=None):
    box(slide, 0, 0, W, Inches(0.12), BLUE)
    box(slide, 0, Inches(0.12), W, Inches(0.92), NAVY)
    add_text(slide, Inches(0.45), Inches(0.22), Inches(12.4), Inches(0.42), title, 24, True, WHITE)
    if subtitle:
        add_text(slide, Inches(0.45), Inches(0.62), Inches(12.4), Inches(0.32), subtitle, 12, False, RGBColor(0xC7, 0xD2, 0xFE))


def bullets(slide, l, t, w, h, items, size=15):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        set_run(p.add_run(), "•  " + item, size, False, DARK)
    return tb


def card(slide, l, t, w, h, title, body, accent=BLUE):
    shape = round_box(slide, l, t, w, h, WHITE)
    shape.line.fill.solid()
    shape.line.color.rgb = LINE
    box(slide, l, t, Inches(0.08), h, accent)
    add_text(slide, l + Inches(0.22), t + Inches(0.12), w - Inches(0.32), Inches(0.32), title, 13, True, NAVY)
    tb = slide.shapes.add_textbox(l + Inches(0.22), t + Inches(0.44), w - Inches(0.36), h - Inches(0.52))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p.add_run(), body, 12, False, MUTED)


# ---------------------------------------------------------------------------
# SLIDE 1 — Title
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, NAVY)
box(s, 0, 0, Inches(0.18), H, BLUE)
add_text(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.35),
         "ACCENTURE INNOVATION CHALLENGE 2026  ·  ROUND 2  ·  TRACK 1", 14, True, RGBColor(0x93, 0xC5, 0xFD))
add_text(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.9), "ControlPlane.ai", 48, True, WHITE)
add_text(s, Inches(0.7), Inches(2.75), Inches(12), Inches(0.5),
         "Enterprise AI Governance & Runtime Risk Control Plane", 22, False, RGBColor(0xC7, 0xD2, 0xFE))
add_text(s, Inches(0.7), Inches(3.5), Inches(11.5), Inches(1.1),
         "Detailed Business Proposal — problem framing, solution design, target users,\nbusiness case and impact, phased roadmap, and key risks with mitigations.",
         16, False, WHITE)
round_box(s, Inches(0.7), Inches(5.15), Inches(4.4), Inches(0.7), BLUE)
add_text(s, Inches(0.7), Inches(5.28), Inches(4.4), Inches(0.5), "Working prototype  ·  Simulated data", 14, True, WHITE, PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(6.15), Inches(11), Inches(0.4),
         "Do not ask “Is this AI response safe?”  Ask: safe for whom, under which policy, with what evidence?",
         14, False, RGBColor(0xBF, 0xDB, 0xFE))

# ---------------------------------------------------------------------------
# SLIDE 2 — Agenda
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "Agenda", "What this proposal covers")
items = [
    ("01", "Problem framing", "Why one generic AI checker fails in a multi-application enterprise"),
    ("02", "Solution design", "Hybrid detectors, policy engine, ALLOW / EDIT / FLAG / BLOCK"),
    ("03", "Target users", "Governance, application owners, reviewers, end users, buyers"),
    ("04", "Business case", "Risk avoidance, review-load reduction, trust, audit readiness"),
    ("05", "Phased roadmap", "Prototype → 90-day pilot → production → agents → platform"),
    ("06", "Risks & mitigations", "Over-flagging, under-flagging, latency, regulation, integration"),
]
for i, (n, title, desc) in enumerate(items):
    col, row = i % 2, i // 2
    l = Inches(0.45 + col * 6.4)
    t = Inches(1.35 + row * 1.8)
    round_box(s, l, t, Inches(6.1), Inches(1.55), WHITE)
    add_text(s, l + Inches(0.25), t + Inches(0.22), Inches(1.0), Inches(0.4), n, 22, True, BLUE)
    add_text(s, l + Inches(1.2), t + Inches(0.28), Inches(4.6), Inches(0.4), title, 18, True, NAVY)
    add_text(s, l + Inches(1.2), t + Inches(0.75), Inches(4.6), Inches(0.55), desc, 13, False, MUTED)
footer(s, 2)

# ---------------------------------------------------------------------------
# SLIDE 3 — Executive summary
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "1. Executive summary", "A control plane, not a one-size-fits-all filter")
add_text(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.7),
         "Enterprises now run many AI systems at once — customer chatbots, employee copilots and decision-support tools. Each has a different risk tolerance. A single generic safety filter either over-flags support chats or under-protects high-impact decisions.",
         15, False, DARK)
card(s, Inches(0.45), Inches(2.05), Inches(12.4), Inches(1.35), "Product",
     "ControlPlane.ai intercepts each AI interaction, evaluates privacy, hallucination, bias and policy risk in parallel, scores the result against the application’s policy, and then allows, edits, flags for human review, or blocks the output. Every decision is explained, evidenced, audited and overridable.",
     BLUE)
add_text(s, Inches(0.5), Inches(3.55), Inches(12.3), Inches(0.7),
         "Core insight: do not ask “Is this AI response safe?” Ask: safe for whom, in which application, under which policy, with what evidence, at what confidence — and what should the enterprise do about it?",
         15, True, NAVY)
for i, (t, b) in enumerate([
    ("What exists now", "Working prototype on 3 simulated use cases, dashboard, review queue, policy simulator, RAG grounding, offline demo mode."),
    ("Ask of this round", "Accept ControlPlane.ai as a credible, demoable enterprise control-plane concept and advance to an enterprise-ready pilot."),
    ("What we will not claim", "No fake production accuracy. No “full GDPR engine.” Simulated data and illustrative planning figures only."),
]):
    card(s, Inches(0.45 + i * 4.2), Inches(4.4), Inches(4.0), Inches(2.5), t, b, BLUE if i < 2 else ORANGE)
footer(s, 3)

# ---------------------------------------------------------------------------
# SLIDE 4 — Problem
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "2. Problem framing", "The enterprise problem is ungoverned AI at scale")
probs = [
    ("Hallucinated facts", "Invented leave, refund, warranty or credit rules enter real workflows."),
    ("Privacy leaks", "Phone, salary, medical or bank details appear in an otherwise ordinary answer."),
    ("Biased reasoning", "Protected attributes used in hiring, lending or customer classification."),
    ("Compounding turns", "A harmless first question becomes a high-risk request three turns later."),
    ("Agentic actions", "When AI can email, update CRM or approve a request, one bad output has real effect."),
    ("Alert fatigue", "A blunt checker flags too much; users bypass it and residual risk returns."),
]
for i, (title, body) in enumerate(probs):
    col, row = i % 3, i // 3
    card(s, Inches(0.4 + col * 4.25), Inches(1.3 + row * 2.7), Inches(4.05), Inches(2.45), title, body,
         RED if i in (1, 2, 4) else ORANGE)
footer(s, 4)

# ---------------------------------------------------------------------------
# SLIDE 5 — One size fails
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "2.1 Why one-size-fits-all fails", "Different use cases have different risk signatures")
rows = [
    ("Use case", "Typical risk", "Latency", "Oversight"),
    ("Customer support chatbot", "PII, wrong policy answers, brand harm", "Tight, real-time", "Moderate"),
    ("Employee knowledge copilot", "Internal data leakage, policy hallucination", "Medium", "Strict"),
    ("Decision-support tool", "Bias, unsupported recommendations, regulation", "Can wait for review", "Very strict"),
]
for r, row in enumerate(rows):
    t = Inches(1.35 + r * 0.7)
    bg = NAVY if r == 0 else WHITE
    fg = WHITE if r == 0 else DARK
    box(s, Inches(0.45), t, Inches(12.4), Inches(0.68), bg)
    widths = [2.8, 4.6, 2.6, 2.4]
    x = Inches(0.55)
    for i, cell in enumerate(row):
        add_text(s, x, t + Inches(0.16), Inches(widths[i]), Inches(0.4), cell, 13, r == 0, fg)
        x += Inches(widths[i])
add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(1.0),
         "A threshold strict enough for credit decisions over-flags ordinary support chats. A threshold loose enough for support under-protects hiring or lending. That is how alert fatigue and liability appear at the same time.",
         16, False, DARK)
round_box(s, Inches(0.45), Inches(5.55), Inches(12.4), Inches(1.3), LIGHT)
add_text(s, Inches(0.7), Inches(5.75), Inches(12.0), Inches(0.9),
         "ControlPlane.ai treats governance as context-aware policy. The same AI output can be FLAG under Customer Support, BLOCK under Employee Copilot, and HUMAN REVIEW under Decision Support.",
         15, True, NAVY)
footer(s, 5)

# ---------------------------------------------------------------------------
# SLIDE 6 — Complexities
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "2.2 Real-world complexities", "The Round 2 brief, designed into the product")
items = [
    ("Overlapping risks", "A fabricated salary for a named employee is both hallucination and a privacy incident. Single-label classification is not enough."),
    ("No real-time ground truth", "The same knowledge gaps that cause hallucination make verification hard. The system must be able to abstain."),
    ("Over- vs under-flagging", "Too many warnings and users bypass the system. Too few and the enterprise is exposed. The tradeoff must be tunable."),
    ("Multi-turn and agents", "“Who is Rahul?” may be harmless. “What is his salary?” and “Give me his phone number” are not."),
    ("Evolving regulation", "Hard-coded legal engines age quickly. Configurable policy by geography and industry is more durable."),
    ("API-based models", "Most enterprises consume foundation models through APIs. The control layer must work at the input/output boundary."),
]
for i, (title, body) in enumerate(items):
    col, row = i % 3, i // 3
    card(s, Inches(0.4 + col * 4.25), Inches(1.3 + row * 2.7), Inches(4.05), Inches(2.45), title, body, BLUE)
footer(s, 6)

# ---------------------------------------------------------------------------
# SLIDE 7 — Assumptions
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "2.3 Operating assumptions", "Directional parameters from the Round 2 brief")
bullets(s, Inches(0.5), Inches(1.3), Inches(7.2), Inches(3.2), [
    "At least three AI applications run at once",
    "Combined volume: tens of thousands of interactions per week",
    "Data sources are a mix of well-governed and loosely governed content",
    "The foundation model is consumed via API",
    "Proprietary production data is not required for this round",
], 16)
round_box(s, Inches(8.0), Inches(1.35), Inches(4.8), Inches(3.4), WHITE)
add_text(s, Inches(8.2), Inches(1.5), Inches(4.4), Inches(0.4), "Illustrative weekly volume", 14, True, NAVY)
for i, (app, n) in enumerate([
    ("Customer Support", "18,000"),
    ("Employee Copilot", "10,000"),
    ("Decision Support", "4,000"),
    ("Total", "32,000"),
]):
    y = Inches(2.05 + i * 0.6)
    add_text(s, Inches(8.3), y, Inches(2.6), Inches(0.4), app, 14, i == 3, DARK)
    add_text(s, Inches(10.8), y, Inches(1.7), Inches(0.4), n, 16, True, BLUE if i < 3 else NAVY)
add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.7),
         "All volumes, interception rates and efficiency figures later in this deck are illustrative planning assumptions. They are not audited customer results and are not claimed as production accuracy.",
         15, False, MUTED)
footer(s, 7)

# ---------------------------------------------------------------------------
# SLIDE 8 — Solution positioning
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "3. Solution design", "Not “an AI that checks another AI”")
round_box(s, Inches(0.45), Inches(1.3), Inches(12.4), Inches(1.5), LIGHT)
add_text(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(1.15),
         "ControlPlane.ai is a policy-aware enterprise control layer that evaluates AI interactions against contextual risk policies and determines whether outputs should be allowed, modified, escalated or blocked. Applications call ControlPlane instead of sending raw model output to the user.",
         16, False, NAVY)
principles = [
    ("Context first", "Same output, different decision by application."),
    ("Hybrid detection", "Rules, retrieval and optional LLM judge, each where they fit."),
    ("Explain every decision", "A score without a reason is not governable."),
    ("Abstain when uncertain", "Low evidence → human review, not false confidence."),
    ("Human in the loop", "High-impact decisions remain reviewable."),
    ("Auditable by default", "Every decision leaves a trail."),
    ("Configurable policy", "Varies by app, region and risk appetite."),
    ("Fail safe", "Detector failure escalates in strict workflows; it does not silently allow."),
]
for i, (title, body) in enumerate(principles):
    col, row = i % 4, i // 4
    card(s, Inches(0.4 + col * 3.2), Inches(3.05 + row * 1.95), Inches(3.05), Inches(1.8), title, body, BLUE)
footer(s, 8)

# ---------------------------------------------------------------------------
# SLIDE 9 — Architecture
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "3.1 Runtime architecture", "Intercept → detect in parallel → score → decide → audit")
steps = [
    ("1", "Pre-gate", "Sensitive requests, injection, restricted topics"),
    ("2", "Generate", "LLM or deterministic mock if input may proceed"),
    ("3", "Detect", "PII, hallucination, bias, policy — in parallel"),
    ("4", "Aggregate", "Weighted score + critical overrides + confidence"),
    ("5", "Decide", "ALLOW / EDIT / FLAG / HUMAN REVIEW / BLOCK"),
    ("6", "Govern", "User output + audit log + review queue + analytics"),
]
for i, (n, title, body) in enumerate(steps):
    col, row = i % 3, i // 3
    l = Inches(0.45 + col * 4.25)
    t = Inches(1.35 + row * 2.55)
    round_box(s, l, t, Inches(4.05), Inches(2.3), WHITE)
    round_box(s, l + Inches(0.2), t + Inches(0.25), Inches(0.5), Inches(0.5), BLUE)
    add_text(s, l + Inches(0.2), t + Inches(0.32), Inches(0.5), Inches(0.4), n, 16, True, WHITE, PP_ALIGN.CENTER)
    add_text(s, l + Inches(0.85), t + Inches(0.3), Inches(3.0), Inches(0.4), title, 18, True, NAVY)
    add_text(s, l + Inches(0.25), t + Inches(1.0), Inches(3.55), Inches(1.0), body, 14, False, MUTED)
footer(s, 9)

# ---------------------------------------------------------------------------
# SLIDE 10 — Detectors
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "3.2 How risk is detected", "Hybrid stack — LLM is never the only source of quantitative truth")
dets = [
    ("Privacy / PII", "Regex and entity detection for phone, email, IDs, bank, salary, medical data. Optional LLM as secondary check. Supports auto-redaction for EDIT."),
    ("Hallucination", "Claim extraction → retrieve enterprise documents → SUPPORTED / CONTRADICTED / UNSUPPORTED / UNVERIFIABLE. Abstain if evidence is missing."),
    ("Bias", "Protected-attribute heuristics plus optional LLM judge. Language is “potential bias.” High-impact cases go to human review."),
    ("Policy", "Enterprise rules: sensitive disclosure, decision oversight, overclaiming, application scope. Example: employee copilot must not make a credit decision."),
]
for i, (title, body) in enumerate(dets):
    col, row = i % 2, i // 2
    card(s, Inches(0.4 + col * 6.4), Inches(1.3 + row * 2.7), Inches(6.2), Inches(2.5), title, body, BLUE)
footer(s, 10)

# ---------------------------------------------------------------------------
# SLIDE 11 — Decision logic
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "3.3 Scoring and decision logic", "Severe individual risk is never hidden by averaging")
add_text(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.55),
         "Overall risk = weighted privacy + hallucination + bias + policy. Weights come from the selected application policy.",
         14, False, DARK)
overrides = [
    ("Privacy ≥ 95", "BLOCK"),
    ("Bias ≥ 90 in very-strict apps", "HUMAN REVIEW"),
    ("Evidence confidence < 25%", "ABSTAIN / REVIEW"),
    ("Detector unavailable", "Fail-safe FLAG / REVIEW"),
]
for i, (k, v) in enumerate(overrides):
    round_box(s, Inches(0.4 + i * 3.2), Inches(1.9), Inches(3.05), Inches(1.25), WHITE)
    add_text(s, Inches(0.5 + i * 3.2), Inches(2.05), Inches(2.85), Inches(0.45), k, 12, True, MUTED)
    add_text(s, Inches(0.5 + i * 3.2), Inches(2.5), Inches(2.85), Inches(0.45), v, 16, True, NAVY)
actions = [
    ("ALLOW", "Low risk. Original response.", GREEN),
    ("EDIT", "Auto-sanitized. Redacted text.", ORANGE),
    ("FLAG", "Meaningful risk. Held for review.", RGBColor(0xEA, 0x58, 0x0C)),
    ("HUMAN REVIEW", "High-impact or uncertain.", VIOLET),
    ("BLOCK", "Critical violation. Safe message.", RED),
]
for i, (title, body, color) in enumerate(actions):
    round_box(s, Inches(0.35 + i * 2.55), Inches(3.45), Inches(2.45), Inches(2.4), WHITE)
    box(s, Inches(0.35 + i * 2.55), Inches(3.45), Inches(2.45), Inches(0.12), color)
    add_text(s, Inches(0.45 + i * 2.55), Inches(3.7), Inches(2.25), Inches(0.7), title, 13, True, color, PP_ALIGN.CENTER)
    add_text(s, Inches(0.45 + i * 2.55), Inches(4.5), Inches(2.25), Inches(1.1), body, 13, False, MUTED, PP_ALIGN.CENTER)
footer(s, 11)

# ---------------------------------------------------------------------------
# SLIDE 12 — Policy proof
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "3.4 Same output. Different policy. Different decision.", "The commercial proof of the control-plane idea")
add_text(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.7),
         "Example response: “The customer request should probably be rejected, although records are incomplete and repayment history could not be verified.”",
         14, False, DARK)
policies = [
    ("Customer Support", "BALANCED", "FLAG", "Can flag uncertain customer handling without freezing the channel.", ORANGE),
    ("Employee Copilot", "STRICT", "BLOCK", "Out of scope: employee copilot must not make a customer credit decision.", RED),
    ("Decision Support", "VERY STRICT", "HUMAN REVIEW", "Right application, but incomplete evidence requires a human.", VIOLET),
]
for i, (app, profile, decision, why, color) in enumerate(policies):
    round_box(s, Inches(0.4 + i * 4.25), Inches(2.1), Inches(4.05), Inches(3.7), WHITE)
    add_text(s, Inches(0.55 + i * 4.25), Inches(2.25), Inches(3.75), Inches(0.4), app, 16, True, NAVY)
    add_text(s, Inches(0.55 + i * 4.25), Inches(2.7), Inches(3.75), Inches(0.35), "Profile: " + profile, 12, False, MUTED)
    round_box(s, Inches(0.7 + i * 4.25), Inches(3.2), Inches(3.45), Inches(0.6), color)
    add_text(s, Inches(0.7 + i * 4.25), Inches(3.3), Inches(3.45), Inches(0.45), decision, 16, True, WHITE, PP_ALIGN.CENTER)
    add_text(s, Inches(0.55 + i * 4.25), Inches(4.05), Inches(3.75), Inches(1.5), why, 14, False, DARK)
footer(s, 12)

# ---------------------------------------------------------------------------
# SLIDE 13 — Users
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "4. Target users", "Sold to the enterprise, used by four roles")
users = [
    ("AI Governance Manager", "Dashboard, Policies, Analytics, Audit Logs", "One place to set policy, prove oversight and produce audit evidence."),
    ("Application Owner", "Analyzer, monitoring, Policy Simulator", "Understand blocks, tune the risk profile, watch false positives and latency."),
    ("Human Reviewer", "Review Queue, interaction detail", "Approve, edit or reject only the uncertain or high-impact cases."),
    ("Enterprise end user", "The original AI app", "Faster trustworthy answers. ControlPlane is invisible unless a response is held or blocked."),
]
for i, (role, surface, need) in enumerate(users):
    card(s, Inches(0.4 + (i % 2) * 6.4), Inches(1.3 + (i // 2) * 2.15), Inches(6.2), Inches(2.0),
         role + "  ·  " + surface, need, BLUE)
add_text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.1),
         "Buying committee: CRO, CISO, Head of AI / CoE, DPO, and line-of-business owners in support, HR, finance or credit. Strongest where AI is already in production and the organisation is regulated or brand-sensitive.",
         14, False, MUTED)
footer(s, 13)

# ---------------------------------------------------------------------------
# SLIDE 14 — Segments
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "4.1 Segments and insertion model", "Do not replace the model. Change the integration point.")
card(s, Inches(0.4), Inches(1.3), Inches(6.2), Inches(3.4), "Near-term segments",
     "Mid-to-large enterprises already using multiple LLM apps. Financial services, insurance, telecom, retail, IT services and shared-services operations. Organisations under GDPR-like, DPDP-like or sector AI-risk pressure.",
     BLUE)
card(s, Inches(6.75), Inches(1.3), Inches(6.15), Inches(3.4), "Later segments",
     "Public sector and healthcare after stronger IAM and data-residency. Multi-model and agentic estates that need tool-call governance, not only text checking. India / EU / US profiles in the prototype are configuration, not legal engines.",
     VIOLET)
round_box(s, Inches(0.4), Inches(4.9), Inches(12.5), Inches(1.85), LIGHT)
add_text(s, Inches(0.65), Inches(5.1), Inches(12.1), Inches(0.4), "Insertion model", 16, True, NAVY)
add_text(s, Inches(0.65), Inches(5.55), Inches(12.1), Inches(0.9),
         "Existing AI app  →  ControlPlane API  →  Foundation model API  →  user-safe output. The client keeps their model vendor. They add a control plane in front of it. Commercial shape: platform subscription by volume and number of apps, plus policy/integration services, plus optional managed review operations.",
         14, False, DARK)
footer(s, 14)

# ---------------------------------------------------------------------------
# SLIDE 15 — Business case
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "5. Business case and impact", "Value is risk-adjusted, not “chatbot conversion +X%”")
values = [
    ("Loss avoidance", "Fewer privacy leaks, fewer ungrounded high-impact recommendations, fewer biased decision traces."),
    ("Operating efficiency", "Humans review only flagged cases, not every AI interaction."),
    ("Trust and adoption", "Business teams scale AI faster when a visible control layer exists."),
    ("Audit readiness", "Every decision has a reason, a policy version and a log."),
    ("Reuse", "One control plane governs many applications instead of a checker per bot."),
]
for i, (title, body) in enumerate(values):
    card(s, Inches(0.35 + i * 2.55), Inches(1.3), Inches(2.45), Inches(3.15), title, body, BLUE)
add_text(s, Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.4),
         "Illustrative mix on 32,000 interactions / week (planning assumption, not a measured customer rate):",
         13, False, MUTED)
for i, (k, v, c) in enumerate([
    ("ALLOW 78%", "Auto-release", GREEN),
    ("EDIT 6%", "Redact and release", ORANGE),
    ("FLAG 11%", "3,520 reviews / week", VIOLET),
    ("BLOCK 5%", "1,600 blocked / week", RED),
]):
    round_box(s, Inches(0.4 + i * 3.2), Inches(5.15), Inches(3.05), Inches(1.6), WHITE)
    add_text(s, Inches(0.5 + i * 3.2), Inches(5.3), Inches(2.85), Inches(0.45), k, 16, True, c)
    add_text(s, Inches(0.5 + i * 3.2), Inches(5.8), Inches(2.85), Inches(0.6), v, 13, False, MUTED)
footer(s, 15)

# ---------------------------------------------------------------------------
# SLIDE 16 — Efficiency
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "5.1 Efficiency and cost of inaction", "Cover more real risk with fewer manual reviews")
round_box(s, Inches(0.4), Inches(1.3), Inches(6.2), Inches(3.55), WHITE)
add_text(s, Inches(0.6), Inches(1.45), Inches(5.8), Inches(0.4), "Without ControlPlane", 16, True, RED)
add_text(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
         "If the organisation spot-checks 20% of conversations:\n\n32,000 × 20% = 6,400 reviews / week\n\nHigh-impact errors can still reach users in the unreviewed 80%. Reviewing everything is not operable.",
         15, False, DARK)
round_box(s, Inches(6.8), Inches(1.3), Inches(6.1), Inches(3.55), WHITE)
add_text(s, Inches(7.0), Inches(1.45), Inches(5.7), Inches(0.4), "With ControlPlane", 16, True, GREEN)
add_text(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(2.5),
         "Review concentrates on the flagged 11%:\n\n32,000 × 11% = 3,520 reviews / week\n\nAbout 2,880 fewer reviews per week, with better coverage of actually risky cases. ALLOW traffic can still be sampled for quality.",
         15, False, DARK)
round_box(s, Inches(0.4), Inches(5.05), Inches(12.5), Inches(1.7), LIGHT)
add_text(s, Inches(0.65), Inches(5.2), Inches(12.1), Inches(0.35), "Cost of inaction is event-driven", 14, True, NAVY)
add_text(s, Inches(0.65), Inches(5.6), Inches(12.1), Inches(0.95),
         "A customer PII leak, a wrong HR/product policy, or biased hiring/credit language can each exceed the cost of a governance layer. Latency is designed as a bounded overhead — if the control plane adds seconds, product teams will route around it.",
         14, False, DARK)
footer(s, 16)

# ---------------------------------------------------------------------------
# SLIDE 17 — Roadmap
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "6. Phased roadmap", "Do not skip from prototype to agent governance")
phases = [
    ("Phase 0", "Now", "Round 2 prototype. Prove intercept, detect, score, decide, explain on simulated data."),
    ("Phase 1", "0–3 months", "Design-partner pilot. Real documents, real model API, SSO/RBAC, measured FP/FN, reviewer SOP."),
    ("Phase 2", "3–9 months", "Enterprise production. Connectors, encryption, retention, FinOps dashboards, shared service."),
    ("Phase 3", "9–18 months", "Agents and multi-model. Intercept tool calls, policy simulation before new use cases."),
    ("Phase 4", "18 months+", "Platform scale. Detector marketplace, sector playbooks, partner-led delivery."),
]
for i, (ph, when, body) in enumerate(phases):
    y = Inches(1.28 + i * 1.1)
    round_box(s, Inches(0.45), y, Inches(12.4), Inches(1.0), WHITE)
    box(s, Inches(0.45), y, Inches(0.1), Inches(1.0), BLUE)
    add_text(s, Inches(0.75), y + Inches(0.12), Inches(2.2), Inches(0.35), ph, 14, True, BLUE)
    add_text(s, Inches(0.75), y + Inches(0.48), Inches(2.2), Inches(0.35), when, 12, False, MUTED)
    add_text(s, Inches(3.2), y + Inches(0.22), Inches(9.3), Inches(0.6), body, 14, False, DARK)
footer(s, 17)

# ---------------------------------------------------------------------------
# SLIDE 18 — Pilot success
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "6.1 90-day pilot success criteria", "What “good” looks like before productising further")
criteria = [
    "At least three applications on-boarded with distinct policies",
    "Every decision is explainable and logged",
    "Reviewers clear the queue within the agreed SLA",
    "False-positive rate is measured and used to tune thresholds",
    "Application owners can show a before/after sample of risky outputs that were edited, flagged or blocked",
    "Added latency stays within the agreed budget for each use case",
]
bullets(s, Inches(0.6), Inches(1.4), Inches(12.0), Inches(4.2), criteria, 18)
add_text(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.8),
         "Recommended sequence: keep the current prototype as the demo artefact → select one design partner with three live AI use cases → run the 90-day policy-tuning pilot → only then productise identity, connectors and agent-action interception.",
         14, False, MUTED)
footer(s, 18)

# ---------------------------------------------------------------------------
# SLIDE 19 — Risks
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, PALE)
header_bar(s, "7. Key risks and mitigations", "The product is designed around these, not despite them")
risks = [
    ("Over-flagging", "Tiered ALLOW / EDIT / FLAG / BLOCK; per-app thresholds; feedback retune; simulator before rollout."),
    ("Under-flagging", "Critical PII overrides; fail-closed for strict apps; human review for bias and low-confidence claims."),
    ("No ground truth", "Retrieval against approved docs; explicit UNVERIFIABLE / abstain; LLM judge is never sole truth."),
    ("Imperfect detectors", "Hybrid stack; “potential bias” language; humans remain in the loop."),
    ("Added latency", "Parallel detectors; cheap prefilters; skip expensive checks where policy allows."),
    ("Policy complexity", "Seeded Balanced / Strict / Very Strict profiles; simulator; split governance vs app owners."),
    ("Regulatory over-claim", "Regional profiles are configuration only; legal review stays human."),
    ("Integration friction", "API-layer insertion; model-agnostic; no need to own or fine-tune the foundation model."),
]
for i, (title, body) in enumerate(risks):
    col, row = i % 4, i // 4
    card(s, Inches(0.3 + col * 3.2), Inches(1.28 + row * 2.7), Inches(3.05), Inches(2.5), title, body, ORANGE)
footer(s, 19)

# ---------------------------------------------------------------------------
# SLIDE 20 — Close
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, W, H, NAVY)
box(s, 0, 0, Inches(0.18), H, BLUE)
add_text(s, Inches(0.7), Inches(1.1), Inches(12), Inches(0.5), "8. Closing recommendation", 14, True, RGBColor(0x93, 0xC5, 0xFD))
add_text(s, Inches(0.7), Inches(1.6), Inches(12), Inches(1.2),
         "Build the control plane before the next wave of AI applications is too large to govern.",
         26, True, WHITE)
add_text(s, Inches(0.7), Inches(3.05), Inches(12), Inches(1.4),
         "Enterprises do not need a perfect oracle. They need a control plane that makes AI use observable, explainable, policy-aware and interruptible.",
         18, False, RGBColor(0xC7, 0xD2, 0xFE))
nexts = [
    "Keep this prototype as the demo and design artefact",
    "Select one design-partner enterprise with three live AI use cases",
    "Replace simulated documents with approved sources",
    "Run a 90-day policy-tuning pilot with measured FP rate, latency and reviewer load",
]
for i, item in enumerate(nexts):
    y = Inches(4.55 + i * 0.4)
    add_text(s, Inches(0.7), y, Inches(12), Inches(0.4), f"{i+1}.  {item}", 14, False, WHITE)
add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
         "ControlPlane.ai  ·  Accenture Innovation Challenge 2026  ·  Round 2  ·  Track 1",
         12, False, RGBColor(0x93, 0xC5, 0xFD))

out = r"C:\Users\Sumit\Music\controlplane ai\docs\ControlPlane_AI_Business_Proposal.pptx"
prs.save(out)
print(out)
print("slides:", len(prs.slides))
